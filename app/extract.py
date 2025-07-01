import os
import pdfplumber
import docx
import pandas as pd
from pptx import Presentation

def extract_text_from_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def extract_text_from_csv_xlsx(file_path):
    try:
        df = pd.read_excel(file_path) if file_path.endswith('.xlsx') else pd.read_csv(file_path)
        return df.astype(str).to_string(index=False)
    except Exception as e:
        raise ValueError(f"Failed to read spreadsheet: {e}")

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".txt":
        return extract_text_from_txt(file_path)
    elif ext in [".csv", ".xlsx"]:
        return extract_text_from_csv_xlsx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")
